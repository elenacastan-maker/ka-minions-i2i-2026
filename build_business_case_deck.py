#!/usr/bin/env python3
"""
Ka-Minions: Business Case Deck Builder
Generates ka-i2i-Business-Case.pptx from the dev-squad template.

3 Value Pillars:
  1. Faster delivery of routine tasks & bugs (3 days → <3 hours)
  2. Zombie task closure — backlog health & noise reduction
  3. On-call support — minimize human participation in guardia shifts

Slide map (after removing empty template slide at index 1):
  0  Cover
  1  The 3 Value Pillars framing   (Problem left / 3 benefits right)
  2  Backlog Data & Evidence
  3  Demo: Ka-Minions in Action
  4  Financial Impact (4 cards)    — pillar by pillar + scale
  5  AI Investment Model (3 cols)  — Paperclip + Claude Code costs vs ROI
  6  4-Phase Roadmap
  7  ANNEX / Risk Mitigation
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
import copy
import lxml.etree as etree

INPUT  = '/Users/elena.castan/ElenaCastan/KA/i2impact/ka-i2i-Minions-dev-squad.pptx'
OUTPUT = '/Users/elena.castan/ElenaCastan/KA/i2impact/ka-i2i-Business-Case.pptx'

prs = Presentation(INPUT)

# ── Helpers ───────────────────────────────────────────────────────────────────

def shape_by_name(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def _make_run(parent_p, text, bold=None, size=None, hex_color=None):
    r_el = etree.SubElement(parent_p, qn('a:r'))
    rPr  = etree.SubElement(r_el, qn('a:rPr'))
    rPr.set('lang', 'en-US')
    if bold is not None:
        rPr.set('b', '1' if bold else '0')
    if size is not None:
        rPr.set('sz', str(int(size * 100)))
    if hex_color:
        fill = etree.SubElement(rPr, qn('a:solidFill'))
        clr  = etree.SubElement(fill, qn('a:srgbClr'))
        clr.set('val', hex_color.lstrip('#'))
    t = etree.SubElement(r_el, qn('a:t'))
    t.text = text


def set_text(shape, text, bold=None, size=None, hex_color=None):
    if shape is None or not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras  = txBody.findall(qn('a:p'))
    template_p = copy.deepcopy(paras[0])
    for p in paras:
        txBody.remove(p)
    new_p = copy.deepcopy(template_p)
    for child in list(new_p):
        tag = etree.QName(child.tag).localname
        if tag in ('r', 'a', 'fld', 'br'):
            new_p.remove(child)
    _make_run(new_p, text, bold=bold, size=size, hex_color=hex_color)
    txBody.append(new_p)


def set_multiline(shape, lines, size=None, hex_color=None, bold=None):
    if shape is None or not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras  = txBody.findall(qn('a:p'))
    template_p = copy.deepcopy(paras[0])
    for p in paras:
        txBody.remove(p)
    for line in lines:
        new_p = copy.deepcopy(template_p)
        for child in list(new_p):
            tag = etree.QName(child.tag).localname
            if tag in ('r', 'a', 'fld', 'br'):
                new_p.remove(child)
        _make_run(new_p, line, bold=bold, size=size, hex_color=hex_color)
        txBody.append(new_p)


BREADCRUMB = 'IDEAS2IMPACT 2026 · Ka-Minions: Business Case'

# ── Delete slide 2 (index 1) — empty template ─────────────────────────────────
sldIdLst = prs.slides._sldIdLst
sldId_to_remove = sldIdLst[1]
prs.part.drop_rel(sldId_to_remove.rId)
sldIdLst.remove(sldId_to_remove)

slides = prs.slides  # 8 slides, indices 0–7


# ── Slide 0: Cover ────────────────────────────────────────────────────────────
s = slides[0]
set_text(shape_by_name(s, 'Google Shape;48;p4'),
         'Ka-Minions: autonomous AI agents that automate routine bugs, '
         'close zombie tasks and reduce on-call burden — €100K+ saved per year, '
         'same headcount')


# ── Slide 1: The 3 Value Pillars ──────────────────────────────────────────────
# Problem left column / 3 benefits right column
s = slides[1]
set_text(shape_by_name(s, 'Google Shape;114;p6'), BREADCRUMB)

# Left column — the pain (warm palette)
set_text(shape_by_name(s, 'Google Shape;122;p6'),
         '30–40% of sprint capacity consumed by bugs, maintenance & ops — not shipping features')
set_text(shape_by_name(s, 'Google Shape;124;p6'),
         '22% of Jira tickets stall >30 days ("zombie tasks"), cluttering backlogs and burying real work')
set_text(shape_by_name(s, 'Google Shape;126;p6'),
         'On-call shifts require manual triage at all hours — high engineer fatigue, slow incident response')

# Right column — the 3 value pillars (green palette)
set_text(shape_by_name(s, 'Google Shape;132;p6'),
         'Pillar 1 — Faster Delivery: resolve ≤1SP bugs & tasks in <3 hours, down from 3+ days')
set_text(shape_by_name(s, 'Google Shape;134;p6'),
         'Pillar 2 — Zombie Task Closure: automated triage, enrichment & resolution of stale backlog debt')
set_text(shape_by_name(s, 'Google Shape;136;p6'),
         'Pillar 3 — On-call Support: Ka-Minions handle Tier-0/1 incidents 24/7, escalating only edge cases to humans')


# ── Slide 2: Backlog Evidence ─────────────────────────────────────────────────
s = slides[2]
set_text(shape_by_name(s, 'Google Shape;152;p7'), BREADCRUMB)
set_text(shape_by_name(s, 'Google Shape;155;p7'), '')
set_text(shape_by_name(s, 'Google Shape;156;p7'),
         'Data source: 6-month Jira backlog — Teams AZ & MO  |  '
         'P1 Faster delivery: avg bug resolution 3 days → target <3 hrs (6× improvement) · '
         'automatable rate 80% of bugs (≤1SP)  |  '
         'P2 Zombie tasks: 22% tickets >30 days · auto-triage + close/re-prioritise via Ka-Minion Centinell mode  |  '
         'P3 On-call: 40% productive resolution + 60% Centinell continuous triage · '
         'reduces P0/P1 human response from hours to minutes')


# ── Slide 3: Demo ─────────────────────────────────────────────────────────────
s = slides[3]
set_text(shape_by_name(s, 'Google Shape;175;p8'), BREADCRUMB)
set_text(shape_by_name(s, 'Google Shape;177;p8'), 'Demo: Ka-Minions in Action')
set_text(shape_by_name(s, 'Google Shape;181;p8'), '')


# ── Slide 4: Financial Impact — pillar by pillar ──────────────────────────────
s = slides[4]
set_text(shape_by_name(s, 'Google Shape;198;p9'), BREADCRUMB)
set_text(shape_by_name(s, 'Google Shape;200;p9'), 'Financial Impact', bold=True, size=29)

# Card 1 — Pillar 1: Faster Delivery (accent: lime #B5E941)
set_text(shape_by_name(s, 'Google Shape;204;p9'), '1. Faster Delivery',      bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;205;p9'),
         '2,281 hrs/yr freed · 80% bugs + 25% BAU tasks automated · '
         '€103,718/yr — 2 pilot teams · payback <1 quarter', size=9)

# Card 2 — Pillar 2: Zombie Tasks (accent: purple #8659E4)
set_text(shape_by_name(s, 'Google Shape;209;p9'), '2. Zombie Task Closure',  bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;210;p9'),
         '22% of backlog stale >30 days · Centinell auto-triages, closes or escalates · '
         'reduces backlog noise by est. 60%, improving sprint planning quality', size=9)

# Card 3 — Pillar 3: On-call (accent: orange #FFA73F)
set_text(shape_by_name(s, 'Google Shape;214;p9'), '3. On-call Support',      bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;215;p9'),
         '24/7 Tier-0/1 incident handling · human escalation only for P0 edge cases · '
         'estimated 30% reduction in out-of-hours engineer interruptions', size=9)

# Card 4 — Phase 4 Scale (accent: teal #2C7C87)
set_text(shape_by_name(s, 'Google Shape;219;p9'), 'Phase 4 — Full KA Scale', bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;220;p9'),
         '80 engineers × 20% efficiency = 16 FTE equivalent · '
         '€1.6M+ value · net gain >€1.2M/yr · 70% team adoption target', size=9)

set_text(shape_by_name(s, 'Google Shape;221;p9'), '')


# ── Slide 5: AI Investment Model — Paperclip + Claude Code ───────────────────
s = slides[5]
set_text(shape_by_name(s, 'Google Shape;233;p10'), BREADCRUMB)
set_text(shape_by_name(s, 'Google Shape;235;p10'),
         'AI Investment Model', bold=True, size=29)

# Column 1 — Paperclip costs
set_multiline(shape_by_name(s, 'Google Shape;236;p10'), [
    'PAPERCLIP COSTS',
    '',
    'Phase 1 setup (one-time):',
    '  Integration + guardrails',
    '  Est. €60K – €80K',
    '',
    'Annual license/maintenance:',
    '  Est. €20K – €30K/yr',
    '',
    'Includes: audit trail,',
    'cost dashboards, kill switch',
    '& human-in-loop gates',
], size=10, hex_color='1D4B00')

# Column 2 — Claude Code / LLM API costs
set_multiline(shape_by_name(s, 'Google Shape;237;p10'), [
    'CLAUDE CODE / LLM COSTS',
    '',
    'Per automated task:',
    '  ~50–100K tokens',
    '  ~€0.80 – €1.50/task',
    '',
    '2,281 tasks/yr (pilot):',
    '  API cost: ~€2K – €3.5K/yr',
    '',
    'Full KA scale (~18K tasks):',
    '  API cost: ~€18K – €27K/yr',
    '',
    'Cost/€ saved ratio: ~3%',
], size=10, hex_color='1D4B00')

# Column 3 — Net ROI summary
set_multiline(shape_by_name(s, 'Google Shape;238;p10'), [
    'NET ROI — SUMMARY',
    '',
    'INVESTMENT (Phase 1):',
    '  Setup:      €60K–80K',
    '  LLM/yr:     €3.5K',
    '  Maint/yr:   €25K',
    '  Total Y1:   ~€90K',
    '',
    'RETURN (Phase 1):',
    '  Savings:    €103,718',
    '',
    'Net gain Y1: >€13K',
    'Break-even:  Q1 pilot',
    'Scale ROI:   3× in Year 3',
], size=10, hex_color='1D4B00')


# ── Slide 6: 4-Phase Roadmap ──────────────────────────────────────────────────
s = slides[6]
set_text(shape_by_name(s, 'Google Shape;254;p11'), BREADCRUMB)

TYPO_MAP = {
    'Ka-Minoins': 'Ka-Minions',
    'Ka-Minois':  'Ka-Minions',
    'Ai-Agents':  'AI Agents',
    'ai-agents':  'AI agents',
    'Howe can':   'How can',
}
for shape in s.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text:
                for wrong, right in TYPO_MAP.items():
                    run.text = run.text.replace(wrong, right)


# ── Slide 7: ANNEX / Risk Mitigation ─────────────────────────────────────────
s = slides[7]
for shape in s.shapes:
    if not shape.has_text_frame:
        continue
    if 'IDEAS2IMPACT 2026' in shape.text_frame.text:
        set_text(shape, BREADCRUMB)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f'Saved → {OUTPUT}')
