#!/usr/bin/env python3
"""
Ka-Minions: C-Level Pitch Deck Builder
Produces a polished 8-slide deck from the hackathon draft.

Slide map after removing empty template slide:
  0  Cover
  1  Problem & The Solution
  2  Use Case Identified with Impact in KA
  3  Demo: Showtime
  4  Solution Architecture          (repurposed from 4-card Impact slide)
  5  Business & Tech Feasibility    (repurposed from rough Business Plan slide)
  6  AI Agents Path to Implementation
  7  ANNEX
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
import copy
import lxml.etree as etree

INPUT  = '/Users/elena.castan/ElenaCastan/KA/i2impact/ka-i2i-Minions-dev-squad.pptx'
OUTPUT = '/Users/elena.castan/ElenaCastan/KA/i2impact/ka-i2i-Minions-C-Level.pptx'

prs = Presentation(INPUT)

# ── Helpers ──────────────────────────────────────────────────────────────────

def shape_by_name(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def _make_run(parent_p, text, bold=None, size=None, hex_color=None):
    """Append a single <a:r> run to a paragraph element."""
    r_el = etree.SubElement(parent_p, qn('a:r'))
    rPr  = etree.SubElement(r_el, qn('a:rPr'))
    rPr.set('lang', 'en-US')
    if bold is not None:
        rPr.set('b', '1' if bold else '0')
    if size is not None:
        rPr.set('sz', str(int(size * 100)))
    if hex_color:
        fill = etree.SubElement(rPr, qn('a:solidFill'))
        clr  = etree.SubElement(fill, qn('a:srgbClr'))
        clr.set('val', hex_color.lstrip('#'))
    t = etree.SubElement(r_el, qn('a:t'))
    t.text = text


def set_text(shape, text, bold=None, size=None, hex_color=None):
    """Replace a shape's entire text frame with a single paragraph + run."""
    if shape is None or not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras  = txBody.findall(qn('a:p'))
    # Save first paragraph as structural template (keeps pPr spacing/align)
    template_p = copy.deepcopy(paras[0])
    for p in paras:
        txBody.remove(p)

    new_p = copy.deepcopy(template_p)
    # Wipe runs/links from the copy
    for child in list(new_p):
        tag = etree.QName(child.tag).localname
        if tag in ('r', 'a', 'fld', 'br'):
            new_p.remove(child)
    _make_run(new_p, text, bold=bold, size=size, hex_color=hex_color)
    txBody.append(new_p)


def set_multiline(shape, lines, size=None, hex_color=None, bold=None):
    """Replace a shape's text frame with one paragraph per line."""
    if shape is None or not shape.has_text_frame:
        return
    txBody = shape.text_frame._txBody
    paras  = txBody.findall(qn('a:p'))
    template_p = copy.deepcopy(paras[0])
    for p in paras:
        txBody.remove(p)

    for line in lines:
        new_p = copy.deepcopy(template_p)
        for child in list(new_p):
            tag = etree.QName(child.tag).localname
            if tag in ('r', 'a', 'fld', 'br'):
                new_p.remove(child)
        _make_run(new_p, line, bold=bold, size=size, hex_color=hex_color)
        txBody.append(new_p)


BREADCRUMB = 'IDEAS2IMPACT 2026 · Ka-Minions: Autonomous AI Dev Agents'

# ── Step 1: Delete slide 2 (index 1) — empty template ────────────────────────
sldIdLst = prs.slides._sldIdLst
sldId_to_remove = sldIdLst[1]
prs.part.drop_rel(sldId_to_remove.rId)
sldIdLst.remove(sldId_to_remove)

slides = prs.slides   # now 8 slides, indices 0–7


# ── Slide 0: Cover ────────────────────────────────────────────────────────────
s = slides[0]
set_text(shape_by_name(s, 'Google Shape;48;p4'),
         'Ka-Minions: autonomous AI agents that tackle tedious tasks and amplify your teams')


# ── Slide 1: Problem & The Solution ──────────────────────────────────────────
s = slides[1]

set_text(shape_by_name(s, 'Google Shape;114;p6'), BREADCRUMB)

# Problem bullets
set_text(shape_by_name(s, 'Google Shape;122;p6'),
         '30–40% of sprint capacity is consumed by maintenance, bug triage & "household chores" — not shipping value')
set_text(shape_by_name(s, 'Google Shape;124;p6'),
         'Routine bug resolution takes >1 working day, stretching release cycles with costly manual iteration loops')
set_text(shape_by_name(s, 'Google Shape;126;p6'),
         'Engineers context-switch constantly on low-value ops, killing focus and slowing product delivery velocity')

# Solution bullets
set_text(shape_by_name(s, 'Google Shape;132;p6'),
         'Ka-Minions: purpose-built AI Agents that own routine dev tasks end-to-end — from detection to resolution')
set_text(shape_by_name(s, 'Google Shape;134;p6'),
         'Governed via Paperclip: centralised cost control, full audit trail & human-in-the-loop approval gates')
set_text(shape_by_name(s, 'Google Shape;136;p6'),
         'Frees 20%+ of engineering time → faster releases, stronger #1 marketplace position, same headcount')


# ── Slide 2: Use Case Identified with Impact in KA ────────────────────────────
s = slides[2]

set_text(shape_by_name(s, 'Google Shape;152;p7'), BREADCRUMB)

# Remove TEMPLATE watermark
set_text(shape_by_name(s, 'Google Shape;155;p7'), '')

# Professional callout over the demo images
set_text(shape_by_name(s, 'Google Shape;156;p7'),
         'Banana Squad pilot: bug auto-triage & resolution in <30 min  ·  PR description generation  ·  '
         'Documentation auto-sync  |  '
         'Why Paperclip? → Governance: transparency, audit trail & centralised agent cost control across all KA teams')


# ── Slide 3: Demo ─────────────────────────────────────────────────────────────
s = slides[3]

set_text(shape_by_name(s, 'Google Shape;175;p8'), BREADCRUMB)          # breadcrumb
set_text(shape_by_name(s, 'Google Shape;177;p8'), 'Demo: Ka-Minions in Action')  # title
# Remove editorial tip
set_text(shape_by_name(s, 'Google Shape;181;p8'), '')


# ── Slide 4: Solution Architecture ───────────────────────────────────────────
# Repurposed from the 4-card Impact & Next Steps template (shapes named *;p9)
s = slides[4]

set_text(shape_by_name(s, 'Google Shape;198;p9'), BREADCRUMB)           # breadcrumb
set_text(shape_by_name(s, 'Google Shape;200;p9'),
         'Solution Architecture', bold=True, size=29)

# Card 1 — User Triggers (accent: lime #B5E941)
set_text(shape_by_name(s, 'Google Shape;204;p9'), 'User Triggers',    bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;205;p9'),
         'Jira ticket · GitHub issue · Slack alert · manual command → any event kicks off a Ka-Minion automatically',
         size=9)

# Card 2 — Ka-Minion Engine (accent: purple #8659E4)
set_text(shape_by_name(s, 'Google Shape;209;p9'), 'Ka-Minion Engine', bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;210;p9'),
         'Paperclip orchestrator + Claude LLM: plan, execute, validate and close tasks — fully autonomously',
         size=9)

# Card 3 — Dev Integrations (accent: orange #FFA73F)
set_text(shape_by_name(s, 'Google Shape;214;p9'), 'Dev Integrations', bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;215;p9'),
         'GitHub API · Jira API · CI/CD pipeline · Confluence docs  →  zero new infrastructure required',
         size=9)

# Card 4 — Governance Layer (accent: teal #2C7C87)
set_text(shape_by_name(s, 'Google Shape;219;p9'), 'Governance Layer', bold=True, size=11)
set_text(shape_by_name(s, 'Google Shape;220;p9'),
         'Human-in-the-loop approval gates · full audit trail · per-agent cost dashboards · instant kill switch',
         size=9)

# Remove editorial note
set_text(shape_by_name(s, 'Google Shape;221;p9'), '')


# ── Slide 5: Business & Tech Feasibility ─────────────────────────────────────
# Repurposed from rough Business Plan slide (shapes named *;p10)
s = slides[5]

set_text(shape_by_name(s, 'Google Shape;233;p10'), BREADCRUMB)
set_text(shape_by_name(s, 'Google Shape;235;p10'),
         'Business & Tech Feasibility', bold=True, size=29)

# Column 1 — Business Impact
set_multiline(shape_by_name(s, 'Google Shape;236;p10'), [
    'BUSINESS IMPACT',
    '',
    '10 engineers  ×  20% reclaimed',
    '= 2 FTE equivalent freed / team',
    '',
    '€500K+ capacity value / year',
    'ROI payback: ~3 months',
    '',
    'Same headcount. More output.',
], size=10, hex_color='1D4B00')

# Column 2 — Tech Stack
set_multiline(shape_by_name(s, 'Google Shape;237;p10'), [
    'TECH STACK',
    '',
    '✓  Paperclip  —  governance layer',
    '✓  Claude API  —  LLM engine',
    '✓  GitHub Actions + Jira webhooks',
    '✓  Existing KA infrastructure',
    '',
    'PoC already running',
    'inside Banana Squad.',
], size=10, hex_color='1D4B00')

# Column 3 — Opportunity Cost
set_multiline(shape_by_name(s, 'Google Shape;238;p10'), [
    'OPPORTUNITY COST',
    '',
    '80 engineers across KA teams',
    '×  20% efficiency gain',
    '= 16 FTE equivalent / year',
    '',
    'Value:       €1.6M+ / year',
    'Build cost:  €400K  / year',
    'Net gain:   >€1.2M  year 1',
], size=10, hex_color='1D4B00')


# ── Slide 6: AI Agents Path to Implementation ─────────────────────────────────
# Fix typos and polish (shapes named *;p11)
s = slides[6]

set_text(shape_by_name(s, 'Google Shape;254;p11'), BREADCRUMB)

TYPO_MAP = {
    'Ka-Minoins': 'Ka-Minions',
    'Ka-Minois':  'Ka-Minions',
    'Ai-Agents':  'AI Agents',
    'ai-agents':  'AI agents',
    'Howe can':   'How can',
}

for shape in s.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text:
                for wrong, right in TYPO_MAP.items():
                    run.text = run.text.replace(wrong, right)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f'Saved → {OUTPUT}')
