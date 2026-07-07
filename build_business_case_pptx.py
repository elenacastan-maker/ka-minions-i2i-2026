#!/usr/bin/env python3
"""
Ka-Minions Business Case — PPTX Builder
Generates a 7-slide deck (6 content + 1 chart) from the financial model.
Run:  python3 build_business_case_pptx.py
Output: ka-minions-business-case.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
import lxml.etree as etree
import copy

OUTPUT = '/Users/elena.castan/ElenaCastan/KA/i2impact/ka-minions-business-case.pptx'

# ── Brand colours ─────────────────────────────────────────────────────────────
G    = RGBColor(0x1D, 0x4B, 0x00)   # dark green
GM   = RGBColor(0x60, 0x9F, 0x28)   # mid green
LIME = RGBColor(0xB5, 0xE9, 0x41)   # lime
PUR  = RGBColor(0x7c, 0x3a, 0xed)   # purple
AMB  = RGBColor(0xd9, 0x77, 0x06)   # amber
TEAL = RGBColor(0x08, 0x91, 0xb2)   # teal
RED  = RGBColor(0xc0, 0x39, 0x2b)   # red
WH   = RGBColor(0xFF, 0xFF, 0xFF)
BK   = RGBColor(0x1a, 0x2e, 0x0d)
GY   = RGBColor(0x50, 0x4E, 0x48)
LGY  = RGBColor(0x6b, 0x7f, 0x5e)

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Financial model data ───────────────────────────────────────────────────────
QUARTERS = ['Q4\'26','Q1\'27','Q2\'27','Q3\'27','Q4\'27','Q1\'28','Q2\'28','Q3\'28','Q4\'28','Q1\'29']
CAL_DATES = ['Q4 2026','Q1 2027','Q2 2027','Q3 2027','Q4 2027','Q1 2028','Q2 2028','Q3 2028','Q4 2028','Q1 2029']

PHASES = [
    {'name':'Phase 1','start':0,'end':0,'teams':2,  'saving_yr':103718, 'color':G},
    {'name':'Phase 2','start':1,'end':2,'teams':7,  'saving_yr':363013, 'color':PUR},
    {'name':'Phase 3','start':3,'end':6,'teams':22, 'saving_yr':1140898,'color':AMB},
    {'name':'Phase 4','start':7,'end':9,'teams':30, 'saving_yr':1557000,'color':TEAL},
]

# Cumulative savings (€K) — 10 quarters, consistent 30t KPI model
CUM_SAVINGS = [25.9, 51.9, 142.6, 233.4, 518.6, 803.8, 1089.0, 1374.3, 1764.3, 2154.3]
# Cumulative investment (€K)
CUM_INVEST  = [87.0, 94.0, 119.0, 129.0, 178.0, 197.0, 216.0, 235.0, 323.0, 361.0]

# Per-phase annual savings for stacked bar chart
PHASE_SAVINGS_YR = {
    'Phase 1': 103.7,   # €K
    'Phase 2': 363.0,
    'Phase 3': 1140.9,
    'Phase 4 (KPI 30t)': 1557.0,
}

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]  # blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=12, bold=False, color=BK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = wrap
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Arial'
    return box


def hdr(slide, title, sub='', eye='IDEAS2IMPACT 2026 · Ka-Minions: Autonomous AI Dev Agents'):
    """Dark green header bar."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill_rgb=G)
    add_rect(slide, 0, 1.1, 13.33, 0.06, fill_rgb=LIME)
    if eye:
        add_text(slide, eye, 0.4, 0.08, 10, 0.25, size=8, bold=True, color=LIME)
    add_text(slide, title, 0.4, 0.3, 10, 0.55, size=18, bold=True, color=WH)
    if sub:
        add_text(slide, sub, 0.4, 0.82, 12.5, 0.28, size=8, color=RGBColor(0xcc,0xff,0x99), italic=True)


def stat_box(slide, label, value, sub, l, t, w=2.4, lime=False):
    bg = RGBColor(0xf0,0xfa,0xe6) if lime else RGBColor(0xf7,0xfa,0xf2)
    bd = LIME if lime else RGBColor(0xD6,0xED,0xAA)
    add_rect(slide, l, t, w, 0.85, fill_rgb=bg, line_rgb=bd, line_pt=1)
    add_text(slide, label, l+0.1, t+0.05, w-0.2, 0.2, size=7, bold=True, color=LGY)
    add_text(slide, value, l+0.1, t+0.22, w-0.2, 0.35, size=20, bold=True, color=G)
    add_text(slide, sub,   l+0.1, t+0.6,  w-0.2, 0.22, size=7.5, color=GY)


def footer_phases(slide):
    """Phase timeline strip at the bottom."""
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill_rgb=RGBColor(0xf7,0xfa,0xf2))
    phases = [
        ('Ph1 — Q4 2026', 'TRX · 2t · €103K/yr', G),
        ('Ph2 — Q1–Q2 2027', '7t · €363K/yr · Break-even Q2 2027', PUR),
        ('Ph3 — Q3 2027–Q2 2028', '22t · €1.14M/yr', AMB),
        ('Ph4 — Q3 2028 → KPI', '30t · €1.56M/yr · >€2.1M net 2028+', TEAL),
    ]
    for i, (title, sub, col) in enumerate(phases):
        x = 0.4 + i * 3.2
        add_rect(slide, x, 7.1, 0.04, 0.4, fill_rgb=col)
        add_text(slide, title, x+0.1, 7.12, 3.0, 0.18, size=7.5, bold=True, color=BK)
        add_text(slide, sub,   x+0.1, 7.28, 3.0, 0.18, size=7,   color=GY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=G)
add_rect(s, 0, 6.8, 13.33, 0.7, fill_rgb=LIME)
add_text(s, 'IDEAS2IMPACT 2026', 0.6, 0.4, 6, 0.4, size=10, bold=True, color=LIME)
add_text(s, 'Ka-Minions', 0.6, 0.9, 8, 1.2, size=52, bold=True, color=LIME)
add_text(s, 'Autonomous AI Dev Agents', 0.6, 2.1, 9, 0.7, size=28, bold=True, color=WH)
add_text(s, 'Free 19.5 FTE equivalents · €1.56M/yr in recovered engineering capacity\n'
            'Same headcount — radically more product delivery',
         0.6, 2.85, 9, 0.7, size=13, color=RGBColor(0xcc,0xff,0x99))
add_text(s, 'Banana Squad · Ideas2Impact 2026 · Adevinta / Kleinanzeigen',
         0.6, 6.85, 8, 0.35, size=9, color=G)
# KPI badges
for i, (v, l) in enumerate([('30/43','Teams KPI target'),('Q2 2027','Break-even'),('~6×','ROI at Q10')]):
    x = 9.8 + i * 1.15
    add_rect(s, x, 1.8, 1.05, 0.7, fill_rgb=RGBColor(0x2a,0x5a,0x10),
             line_rgb=LIME, line_pt=1)
    add_text(s, v, x+0.05, 1.85, 0.95, 0.35, size=16, bold=True, color=LIME, align=PP_ALIGN.CENTER)
    add_text(s, l, x+0.05, 2.2,  0.95, 0.25, size=7,  color=WH, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & KPIs
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
hdr(s, 'The Problem: 30–40% of engineering capacity is lost to maintenance & ops',
    sub='KPI target: 70% of KA teams (30/43) with active AI agents · Full potential at 43t: €2.23M/yr')

# Stat tiles row
stats = [
    ('Avg bug resolution', '3+ days', '→ <3 hours with Ka-Minions', False),
    ('Zombie tickets', '22%', '>30 days stale across all teams → <8%', False),
    ('Out-of-hours calls & alerts', '3/week', '→ −30% with Sentinel coverage', False),
    ('Sprint on ops/maintenance', '30–40%', '→ +20% capacity freed', False),
    ('Recoverable capacity (30t)', '€1.56M/yr', 'from Q3 2028 · 19.5 FTE equivalent', True),
]
for i, (lbl, val, sub, lime) in enumerate(stats):
    stat_box(s, lbl, val, sub, 0.3 + i*2.55, 1.28, w=2.4, lime=lime)

# Three pillars
pillar_data = [
    ('Pillar 01', 'Faster bug resolution\n& routine tasks',
     'BEFORE: 3+ days average, blocks releases\nKa-Minion: detects → fixes → merges in <3 hours\nPhase 1 gate KPI: 3d → <3h'),
    ('Pillar 02', 'Zombie ticket\nelimination',
     'BEFORE: 22% of tickets stale >30 days\nKa-Minion: auto-closes confirmed zombies\nPhase 2 KPI: <8% stale rate'),
    ('Pillar 03', 'Out-of-hours duty\ncalls & alerts',
     'BEFORE: ~3 hands-on calls & alerts/week/team\nKa-Minion Sentinels: watch 24/7, silence false alerts\nPhase 2 KPI: −30% hands-on duty calls'),
    ('Pillar 04', 'Sprint capacity\nreclaimed',
     'BEFORE: 30–40% sprint on ops housekeeping\nKa-Minions own BAU tasks end-to-end\nFreed: +20% capacity for product delivery'),
]
for i, (pn, title, body) in enumerate(pillar_data):
    x = 0.3 + i * 3.18
    add_rect(s, x, 2.32, 3.0, 4.58, fill_rgb=RGBColor(0xf7,0xfa,0xf2),
             line_rgb=RGBColor(0xD6,0xED,0xAA), line_pt=1)
    add_rect(s, x, 2.32, 3.0, 0.04, fill_rgb=[G,PUR,AMB,TEAL][i])
    add_text(s, pn,    x+0.15, 2.38, 2.7, 0.22, size=7.5, bold=True, color=LGY)
    add_text(s, title, x+0.15, 2.58, 2.7, 0.45, size=12,  bold=True, color=G)
    add_text(s, body,  x+0.15, 3.1,  2.7, 3.5,  size=9.5, color=BK)

footer_phases(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Team Expansion & ROI by Domain
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
hdr(s, '€1.56M/yr at 70% adoption — phased rollout across 9 KA domains',
    sub='Phase 1: TRX pilot (2 teams, €103K/yr) · Phase 4 KPI: 30/43 teams (70%) · Full 43t: €2.23M')

# Phase investment table
add_text(s, 'Phase Investment & Savings Summary', 0.4, 1.25, 6, 0.3, size=10, bold=True, color=G)
table_data = [
    ['Phase', 'Period', 'Teams', 'Annual Savings', 'Running Cost/Q', 'Net/Q'],
    ['Phase 1', 'Q4 2026', '2', '€103,718', '~€3,500', '~€22K'],
    ['Phase 2', 'Q1–Q2 2027', '7', '€363,013', '~€9,000', '~€82K'],
    ['Phase 3', 'Q3 2027–Q2 2028', '22', '€1,140,898', '~€25,000', '~€260K'],
    ['Phase 4 KPI (30t)', 'Q3 2028+', '30', '€1,557,000', '~€80K/yr', '~€370K'],
    ['Full potential (43t)', '2028+', '43', '€2,229,937', '~€115K/yr', '~€443K'],
]
tbl = s.shapes.add_table(6, 6, Inches(0.4), Inches(1.6), Inches(7.5), Inches(2.0)).table
col_widths = [1.5, 1.6, 0.7, 1.5, 1.5, 0.7]
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = Inches(cw)
for ri, row in enumerate(table_data):
    for ci, cell_text in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci >= 2 else PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(9 if ri > 0 else 8)
        run.font.bold = (ri == 0 or ri == 5)
        run.font.color.rgb = WH if ri == 0 else (G if ri == 4 else BK)
        run.font.name = 'Arial'
        fill = cell.fill
        if ri == 0:
            fill.solid(); fill.fore_color.rgb = G
        elif ri == 4:
            fill.solid(); fill.fore_color.rgb = RGBColor(0xf0,0xfa,0xe6)
        elif ri % 2 == 0:
            fill.solid(); fill.fore_color.rgb = RGBColor(0xfa,0xfd,0xf5)

# ROI metrics
roi_tiles = [
    ('Total investment', '€361K', 'Q4 2026 → Q10'),
    ('Break-even', 'Q2 2027', 'savings overtake cost'),
    ('ROI at Q10', '~6×', '€2.15M saved vs €361K'),
    ('Net value 2028+', '>€2.1M', '30t KPI · 19.5 FTE freed'),
    ('AI stack cost', '~3%', 'of savings (Paperclip + Claude)'),
]
for i, (lbl, val, sub) in enumerate(roi_tiles):
    x = 8.2 + (i % 3) * 1.68
    y = 1.55 + (i // 3) * 1.05
    lime = (i == 3)
    stat_box(s, lbl, val, sub, x, y, w=1.6, lime=lime)

footer_phases(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Costs, ROI & Governance
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
hdr(s, 'Enabler team + two SaaS tools — total ~€115K/yr vs €1.56M/yr freed',
    sub='1 EM + 2 engineers (existing headcount reallocation) · Paperclip + Claude API · cost ≈ 7% of KPI savings')

sections = [
    ('Enabler Team', G, [
        'EM (50% Phase 1 → full from Ph2)',
        '2 Engineers (1 from Ph1, +1 from Ph2)',
        'Prioritises backlog, manages phase gates',
        'Trains domain-specific agent context',
        '→ No new hire needed until Phase 2',
    ]),
    ('Paperclip Platform', PUR, [
        'Orchestration & governance layer',
        '~€12–24K/yr (team plan)',
        'Full audit trail per agent action',
        'Human-in-the-loop approval gates',
        'Instant kill-switch per agent or org-wide',
    ]),
    ('Claude API', AMB, [
        'LLM engine for agent reasoning',
        '~€1/task average across all types',
        'Phase 1 (3.5K tasks): ~€3.5K/yr',
        'Phase 4 (52K tasks): ~€52K/yr',
        'Cost ≈ 3.3% of savings at scale',
    ]),
    ('Governance & Risk', TEAL, [
        'Scope guardrail: ≤1 SP tasks only',
        'Budget cap per agent (EM override)',
        'CI must be green before any merge',
        'Per-agent cost dashboard in Paperclip',
        'Rollback in 1 click — zero blast radius',
    ]),
]

for i, (title, col, bullets) in enumerate(sections):
    x = 0.3 + i * 3.22
    add_rect(s, x, 1.25, 3.1, 5.65, fill_rgb=RGBColor(0xf9,0xfc,0xf5),
             line_rgb=RGBColor(0xD6,0xED,0xAA), line_pt=1)
    add_rect(s, x, 1.25, 3.1, 0.05, fill_rgb=col)
    add_text(s, title, x+0.15, 1.32, 2.8, 0.3, size=11, bold=True, color=col)
    for j, b in enumerate(bullets):
        add_text(s, '›  ' + b, x+0.12, 1.72 + j*0.75, 2.85, 0.68, size=9.5, color=BK)

footer_phases(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Team & AI Cost Model
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
hdr(s, 'Team & AI Cost Model — phase-by-phase breakdown',
    sub='All costs include one-time setup, LLM API, maintenance · LLM ≈ 3% of savings · governance via Paperclip from day one')

# Team cost table
add_text(s, 'Enabler team headcount & cost by phase', 0.4, 1.28, 6.5, 0.28, size=10, bold=True, color=G)
team_rows = [
    ['Phase', 'Period', 'Team', 'People Cost', 'Tech Stack', 'Total Cost', 'Quarterly Savings', 'Net/Q'],
    ['Phase 1', 'Q4 2026', '0.5 EM + 1 eng', '~€40K/yr', '~€3.5K/yr', '~€74K setup+run', '€25.9K', '−€48K (invest)'],
    ['Phase 2', 'Q1–Q2 2027', '1 EM + 2 eng', '~€150K/yr', '~€15K/yr', '~€165K/yr', '€90.8K/Q', '+€67K/Q net'],
    ['Phase 3', 'Q3 2027–Q2 2028', '1 EM + 2 eng', '~€150K/yr', '~€50K/yr', '~€200K/yr', '€285.2K/Q', '+€235K/Q net'],
    ['Phase 4 (KPI 30t)', 'Q3 2028+', '1 EM + 2 eng', '~€150K/yr', '~€80K/yr', '~€230K/yr', '€390K/Q', '+€332K/Q net'],
]
tbl2 = s.shapes.add_table(5, 8, Inches(0.4), Inches(1.62), Inches(12.5), Inches(1.85)).table
cws = [0.85, 1.5, 1.3, 1.1, 1.0, 1.3, 1.4, 1.55]
for ci, cw in enumerate(cws):
    tbl2.columns[ci].width = Inches(cw)
for ri, row in enumerate(team_rows):
    for ci, txt in enumerate(row):
        cell = tbl2.cell(ri, ci)
        cell.text = txt
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci >= 3 else PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(8 if ri > 0 else 7.5)
        run.font.bold = (ri == 0)
        run.font.color.rgb = WH if ri == 0 else (G if (ri == 4 and ci >= 6) else BK)
        run.font.name = 'Arial'
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = G
        elif ri % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xfa,0xfd,0xf5)

# AI cost model
add_text(s, 'AI Stack cost model (Paperclip + Claude API)', 0.4, 3.7, 8, 0.28, size=10, bold=True, color=G)
ai_data = [
    ('Phase 1 · 2 teams · ~3,500 tasks/yr', '€3,500', '~3.4% of €103K savings', AMB),
    ('Phase 2 · 7 teams · ~12,250 tasks/yr', '€12,250', '~3.4% of €363K savings', PUR),
    ('Phase 3 · 22 teams · ~38,500 tasks/yr', '€38,500', '~3.4% of €1.14M savings', AMB),
    ('Phase 4 KPI · 30 teams · ~52,500 tasks/yr', '€52,500', '~3.4% of €1.56M savings', TEAL),
]
for i, (lbl, cost, pct, col) in enumerate(ai_data):
    x = 0.4 + i * 3.2
    add_rect(s, x, 4.05, 3.05, 0.75, fill_rgb=RGBColor(0xf7,0xfa,0xf2),
             line_rgb=RGBColor(0xD6,0xED,0xAA), line_pt=1)
    add_rect(s, x, 4.05, 3.05, 0.04, fill_rgb=col)
    add_text(s, lbl,  x+0.12, 4.12, 2.8, 0.25, size=8,   color=GY)
    add_text(s, cost, x+0.12, 4.3,  1.2, 0.3,  size=14,  bold=True, color=col)
    add_text(s, pct,  x+1.4,  4.38, 1.5, 0.22, size=8,   color=LGY)

# Insight
add_rect(s, 0.4, 5.0, 12.5, 0.5, fill_rgb=RGBColor(0xf0,0xfa,0xe6),
         line_rgb=RGBColor(0xB5,0xE9,0x41), line_pt=1)
add_text(s, '💡  Phase-gated: no spend unlocks without hitting the prior gate KPI. '
            'Single ask today: €60–80K Phase 1 setup (Q4 2026). Break-even Q2 2027. €1.56M/yr from Q3 2028. ROI >6× by 2028+.',
         0.55, 5.07, 12.2, 0.38, size=9.5, color=G)

footer_phases(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Impact at a Glance (C-Level) with Savings Ramp Chart
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WH)
hdr(s, 'Autonomous AI agents that free engineering capacity worth €1.56M/yr —\nsame headcount, radically more product delivery',
    sub='KPI target: 30/43 teams (70%) · Break-even Q2 2027 · ROI ~6× at Q10 · >€2.1M net in 2028+')

# Badge row
for i, (v, l) in enumerate([('30/43','Teams KPI'), ('Q2 2027','Break-even'), ('~6×','ROI Q10'), ('>€2.1M','Net 2028+')]):
    x = 9.2 + i * 1.02
    add_rect(s, x, 0.15, 0.95, 0.65, fill_rgb=RGBColor(0x2a,0x5a,0x10),
             line_rgb=LIME, line_pt=1)
    add_text(s, v, x+0.05, 0.17, 0.85, 0.3, size=13, bold=True, color=LIME, align=PP_ALIGN.CENTER)
    add_text(s, l, x+0.05, 0.46, 0.85, 0.2, size=6.5, color=WH, align=PP_ALIGN.CENTER)

# COL 1 — The Problem
add_text(s, 'THE PROBLEM TODAY', 0.35, 1.22, 4.0, 0.22, size=8, bold=True, color=GM)
add_text(s, 'Engineers doing work AI agents can own', 0.35, 1.42, 4.0, 0.35, size=11, bold=True, color=G)
pairs = [
    ('3+ days','<3 h','Bug resolution · 6× faster','Phase 1 gate KPI'),
    ('22%','<8%','Zombie tickets · stale >30d','Phase 2 KPI'),
    ('3/wk','−30%','Out-of-hours calls & alerts','Phase 2 KPI'),
    ('30–40%','+20%','Sprint on ops → capacity freed','Freed'),
]
for i, (bef, aft, sub, kpi) in enumerate(pairs):
    y = 1.85 + i * 1.18
    add_rect(s, 0.35, y, 4.0, 1.08, fill_rgb=RGBColor(0xff,0xf5,0xf5))
    add_text(s, bef, 0.45, y+0.08, 0.9, 0.45, size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(s, '→',  1.35, y+0.18, 0.4, 0.3, size=16, color=GM, align=PP_ALIGN.CENTER)
    add_text(s, aft,  1.75, y+0.08, 1.0, 0.45, size=18, bold=True, color=G, align=PP_ALIGN.CENTER)
    add_text(s, kpi,  2.8, y+0.22, 1.45, 0.22, size=7.5, bold=True, color=GY, align=PP_ALIGN.CENTER)
    add_text(s, sub,  0.38, y+0.76, 3.9, 0.25, size=8.5, color=GY)

# Divider
add_rect(s, 4.45, 1.18, 0.02, 5.72, fill_rgb=RGBColor(0xe5,0xf0,0xd2))

# COL 2 — Financial Chart (Savings Ramp)
add_text(s, 'FINANCIAL IMPACT · Q4 2026 → 2028+', 4.6, 1.22, 4.5, 0.22, size=8, bold=True, color=GM)
add_text(s, 'Same headcount. Radically more output.', 4.6, 1.42, 4.5, 0.3, size=11, bold=True, color=G)

# KPI tiles
for i, (lbl, val, sub, lime) in enumerate([
    ('Savings at KPI target', '€1.56M/yr', '30 teams · from Q3 2028', True),
    ('FTE capacity freed', '19.5 eng', 'same headcount', True),
    ('ROI at Q10', '~6×', '€2.15M vs €361K invested', False),
]):
    x = 4.6 + i * 1.58
    stat_box(s, lbl, val, sub, x, 1.8, w=1.5, lime=lime)

# Savings ramp bar chart
chart_data = ChartData()
chart_data.categories = ['Q4\'26\n€103K', 'Q1\'27\n€363K', 'Q3\'27\n€1.14M', 'Q3\'28\n€1.56M', '2028+\n>€2.1M']
chart_data.add_series('Phase 1 (2t)', (103.7,    0,      0,      0,      0))
chart_data.add_series('Phase 2 (7t)', (0,       363.0,   0,      0,      0))
chart_data.add_series('Phase 3 (22t)',(0,       0,     1140.9,   0,      0))
chart_data.add_series('Phase 4 KPI', (0,        0,      0,     1557.0,   0))
chart_data.add_series('2028+ full',  (0,        0,      0,      0,     2100.0))

chart = s.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(4.6), Inches(2.78), Inches(4.35), Inches(3.05),
    chart_data
).chart

chart.chart_style = 2
from pptx.enum.chart import XL_LEGEND_POSITION
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

series_colors = [G, PUR, AMB, TEAL, LIME]
for i, series in enumerate(chart.series):
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = series_colors[i]

# Divider
add_rect(s, 9.05, 1.18, 0.02, 5.72, fill_rgb=RGBColor(0xe5,0xf0,0xd2))

# COL 3 — How it works
add_text(s, 'HOW Ka-Minions WORK', 9.18, 1.22, 3.9, 0.22, size=8, bold=True, color=GM)
add_text(s, 'Jira ticket → merged PR —\nautonomous, governed, audited', 9.18, 1.42, 3.9, 0.45, size=10, bold=True, color=G)

flow = [
    ('📥', G,    'Trigger', 'Jira · GitHub · Slack · manual\n→ any ≤1SP task auto-detected'),
    ('🧠', PUR,  'Paperclip + Claude', 'Reads context, writes fix,\nruns CI · 24/7 · no queue'),
    ('✅', TEAL, 'Approve & Merge', 'CI green + Paperclip audit gate\nHuman-in-loop · full audit trail'),
]
for i, (icon, col, title, body) in enumerate(flow):
    y = 1.95 + i * 1.32
    add_rect(s, 9.18, y, 3.9, 1.2, fill_rgb=RGBColor(0xf6,0xfa,0xf0),
             line_rgb=RGBColor(0xD6,0xED,0xAA), line_pt=1)
    add_text(s, icon,  9.22, y+0.05, 0.35, 0.35, size=16)
    add_text(s, title, 9.6,  y+0.08, 3.4, 0.28, size=10, bold=True, color=col)
    add_text(s, body,  9.6,  y+0.38, 3.35, 0.55, size=9, color=BK)

# Governance grid
add_text(s, 'GOVERNANCE GUARDRAILS', 9.18, 5.94, 3.9, 0.22, size=8, bold=True, color=GM)
gov = ['🔒 Only ≤1SP scope','👁️ Full Paperclip audit trail','🛑 Instant kill-switch','💰 Per-agent budget caps']
for i, g in enumerate(gov):
    x = 9.18 + (i % 2) * 1.95
    y = 6.2 + (i // 2) * 0.42
    add_rect(s, x, y, 1.88, 0.38, fill_rgb=RGBColor(0xf0,0xfa,0xe6),
             line_rgb=RGBColor(0xD6,0xED,0xAA), line_pt=1)
    add_text(s, g, x+0.08, y+0.07, 1.72, 0.26, size=8, color=G)

footer_phases(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Path to Implementation
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0xF4,0xF2,0xEF))
add_text(s, 'IDEAS2IMPACT 2026 · ka-i2i-autonomous-minions-agents',
         0.4, 0.12, 10, 0.22, size=8, bold=True, color=GM)
add_text(s, 'Ka-Minions path to implementation',
         0.4, 0.35, 12, 0.6, size=28, bold=True, color=G)

# NOW/NEXT banner
add_rect(s, 0.35, 1.05, 6.1,  0.45, fill_rgb=PUR)
add_rect(s, 6.55, 1.05, 6.43, 0.45, fill_rgb=GM)
add_text(s, 'NOW: Training & learning',         0.5, 1.1, 5.8, 0.35, size=13, bold=True, color=WH)
add_text(s, 'NEXT: Drive business value delivery', 6.65, 1.1, 6.1, 0.35, size=13, bold=True, color=WH)

phase_cards = [
    ('PHASE 1\nHypothesis validation', 'Q4 2026 · 1Q',
     'Automate repetitive\noperations',
     'Pilot teams (1–2) creating and\ntraining their own Ka-Minions',
     '€103K/yr · 2 teams',
     'KPI: Bug resolution 3 days → <3 hours',
     '1Q', G),
    ('PHASE 2\nMVPs creation', 'Q1–Q2 2027 · 2Q',
     'Boosting KA Impact:\nfuel teams with AI Agents',
     'Five selected teams implementing\nAI agents for product delivery',
     '€363K/yr · 7 teams\n✓ Break-even Q2 2027',
     'KPI: 70% success rate Ka-Minions usage',
     '2 Qs', PUR),
    ('PHASE 3\nGrowth & Iterations', 'Q3 2027–Q2 2028 · 4Q',
     'AI-Agents\nAcceleration',
     'Driving efficiency by integrating\nAI agents within existing headcounts',
     '€1.14M/yr · 22 teams',
     'KPI: 40% Teams Ka-Minions adoption',
     '4 Qs', AMB),
    ('PHASE 4\nScale Up', 'Q3 2028 → Always on',
     'AI agentic teams',
     'Building Features and Services\nwith AI Agents as Team Members',
     '€1.56M/yr · 30 teams · KPI target',
     'KPI: 70% adoption + 30% feature delivery',
     'Always on', TEAL),
]

for i, (header, period, title, desc, saving, kpi, dur, col) in enumerate(phase_cards):
    x = 0.35 + i * 3.22
    add_rect(s, x, 1.6, 3.1, 5.15, fill_rgb=WH, line_rgb=RGBColor(0xD6,0xED,0xAA), line_pt=1)
    add_rect(s, x, 1.6, 3.1, 0.05, fill_rgb=col)
    bg_h = RGBColor(0xf0,0xfa,0xe6) if col==G else \
           RGBColor(0xf5,0xf0,0xff) if col==PUR else \
           RGBColor(0xff,0xfb,0xeb) if col==AMB else RGBColor(0xf0,0xf9,0xff)
    add_rect(s, x, 1.65, 3.1, 0.55, fill_rgb=bg_h)
    add_text(s, header, x+0.12, 1.67, 2.85, 0.3, size=8, bold=True, color=col)
    add_text(s, period, x+0.12, 1.95, 2.85, 0.2, size=7.5, color=GY)
    add_text(s, title,  x+0.12, 2.28, 2.85, 0.45, size=11, bold=True, color=col)
    add_text(s, desc,   x+0.12, 2.78, 2.85, 0.6,  size=9,  color=GY, italic=True)

    # Saving chip
    add_rect(s, x+0.1, 3.48, 2.9, 0.42, fill_rgb=bg_h, line_rgb=col, line_pt=1)
    add_text(s, saving, x+0.2, 3.54, 2.7, 0.3, size=9.5, bold=True, color=col)

    # KPI box
    add_rect(s, x+0.1, 4.0, 2.9, 0.4, fill_rgb=RGBColor(0xf0,0xfa,0xe6))
    add_text(s, kpi, x+0.18, 4.06, 2.7, 0.28, size=8.5, bold=True, color=G)

    # Duration
    add_text(s, dur, x+1.3, 5.88, 0.8, 0.35, size=17, bold=True, color=col, align=PP_ALIGN.CENTER)

# ROI footer bar
add_rect(s, 0, 6.78, 13.33, 0.72, fill_rgb=G)
add_rect(s, 0, 6.78, 13.33, 0.04, fill_rgb=LIME)
roi_footer = [
    ('Total investment', '€361K', 'Q4 2026 → Q10'),
    ('Break-even', 'Q2 2027', 'savings overtake cost'),
    ('Savings at KPI', '€1.56M/yr', '30 teams · from Q3 2028'),
    ('ROI at Q10', '~6×', '€2.15M vs €361K invested'),
    ('Net value 2028+', '>€2.1M', 'full potential 43 teams'),
]
for i, (lbl, val, sub) in enumerate(roi_footer):
    x = 0.5 + i * 2.6
    add_text(s, lbl, x, 6.82, 2.4, 0.2, size=7, bold=True,
             color=RGBColor(0xcc,0xcc,0xcc), align=PP_ALIGN.CENTER)
    add_text(s, val, x, 6.98, 2.4, 0.3, size=14, bold=True,
             color=LIME, align=PP_ALIGN.CENTER)
    add_text(s, sub, x, 7.25, 2.4, 0.18, size=6.5,
             color=RGBColor(0x99,0xbb,0x99), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# CHART SLIDE — Investment vs Cumulative Savings (10 quarters)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WH)
hdr(s, 'Cumulative savings overtake investment at Q2 2027 — ~€2.15M saved (30t KPI) vs €361K by Q10',
    sub='Phase-gated · 30/43 teams (70% KPI) · LLM cost ≈ 3% of savings · Savings valued at €80K/yr avg engineer · Break-even Q2 2027')

# Line chart: investment vs savings
cd = ChartData()
cd.categories = CAL_DATES
cd.add_series('Cumulative Savings (€K, 30t KPI)',
    (25.9, 51.9, 142.6, 233.4, 518.6, 803.8, 1089.0, 1374.3, 1764.3, 2154.3))
cd.add_series('Cumulative Investment (€K)',
    (87.0, 94.0, 119.0, 129.0, 178.0, 197.0, 216.0, 235.0, 323.0, 361.0))

ch = s.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(0.5), Inches(1.3), Inches(9.0), Inches(5.7),
    cd
).chart

ch.chart_style = 10
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.BOTTOM

series_clrs = [G, RED]
for i, ser in enumerate(ch.series):
    ser.format.line.color.rgb = series_clrs[i]
    ser.format.line.width = Pt(2.5 if i == 0 else 1.5)

# Annotations panel
add_rect(s, 9.7, 1.3, 3.3, 5.7, fill_rgb=RGBColor(0xf7,0xfa,0xf2),
         line_rgb=RGBColor(0xD6,0xED,0xAA), line_pt=1)
add_text(s, 'KEY MILESTONES', 9.9, 1.4, 2.9, 0.25, size=8, bold=True, color=GM)
milestones = [
    (G,   'Q4 2026',   'Phase 1 launch',     '2 teams · €103K/yr\n€87K setup investment'),
    (PUR, 'Q1 2027',   'Phase 2 expansion',  '7 teams · €363K/yr'),
    (AMB, '✓ Q2 2027', 'BREAK-EVEN',         'Savings overtake cost'),
    (AMB, 'Q3 2027',   'Phase 3 ramp',       '22 teams · €1.14M/yr\n+15 teams added'),
    (TEAL,'Q3 2028',   'Phase 4 KPI target', '30 teams · €1.56M/yr\n19.5 FTE freed'),
    (LIME,'2028+',     'Full potential',      '>€2.1M net per year'),
]
for i, (col, date, title, body) in enumerate(milestones):
    y = 1.8 + i * 0.82
    add_rect(s, 9.78, y, 0.04, 0.55, fill_rgb=col)
    add_text(s, date,  9.88, y,      2.9, 0.22, size=8.5, bold=True, color=col)
    add_text(s, title, 9.88, y+0.2,  2.9, 0.2,  size=9,   bold=True, color=BK)
    add_text(s, body,  9.88, y+0.38, 2.9, 0.35, size=8,   color=GY)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f'✅  Saved → {OUTPUT}')
